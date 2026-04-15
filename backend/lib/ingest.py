"""
IRIS ingestion layer.
Parses uploaded files into a single plain-text document for Claude extraction.

Supported formats:
  TXT, CSV, EML, DOCX, PDF, PPTX, XLSX, HTML, JSON
  Images (PNG, JPG, JPEG, GIF, WEBP) — handled via Claude vision in server.py
  ZIP — extracted and each file parsed by extension
"""

import csv
import io
import email
from email import policy


def parse_text(content: str) -> str:
    """Return raw text as-is."""
    return content.strip()


def parse_csv(content: str, filename: str = '') -> str:
    """Flatten a CSV into a labelled text block."""
    reader = csv.DictReader(io.StringIO(content))
    rows = list(reader)
    if not rows:
        return ''

    lines = [f'[CSV file: {filename}]' if filename else '[CSV data]']
    lines.append(f'Columns: {", ".join(rows[0].keys())}')
    lines.append(f'Rows: {len(rows)}')
    lines.append('')

    for i, row in enumerate(rows, 1):
        parts = [f'{k}: {v}' for k, v in row.items() if v and v.strip()]
        lines.append(f'Row {i}: {" | ".join(parts)}')

    return '\n'.join(lines)


def parse_eml(content: bytes | str, filename: str = '') -> str:
    """Extract subject, sender, date, and body from an EML file."""
    if isinstance(content, str):
        content = content.encode('utf-8', errors='replace')

    msg = email.message_from_bytes(content, policy=policy.default)
    lines = [f'[Email file: {filename}]' if filename else '[Email]']

    for header in ('From', 'To', 'Subject', 'Date'):
        val = msg.get(header, '')
        if val:
            lines.append(f'{header}: {val}')

    lines.append('')

    body = ''
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_type() == 'text/plain':
                body = part.get_content()
                break
    else:
        if msg.get_content_type() == 'text/plain':
            body = msg.get_content()

    if body:
        lines.append('Body:')
        lines.append(body.strip())

    return '\n'.join(lines)


def parse_docx(content: bytes, filename: str = '') -> str:
    """Extract all text from a .docx file, preserving heading structure."""
    from docx import Document

    doc = Document(io.BytesIO(content))
    lines = [f'[Word document: {filename}]' if filename else '[Word document]']
    lines.append('')

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ''
        if style.startswith('Heading'):
            level = style.replace('Heading', '').strip()
            prefix = '#' * int(level) if level.isdigit() else '#'
            lines.append(f'\n{prefix} {text}')
        else:
            lines.append(text)

    for table in doc.tables:
        lines.append('')
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                lines.append(' | '.join(cells))

    return '\n'.join(lines)


def parse_pdf(content: bytes, filename: str = '') -> str:
    """Extract text from a PDF using pdfplumber."""
    import pdfplumber

    lines = [f'[PDF file: {filename}]' if filename else '[PDF document]']
    lines.append('')

    with pdfplumber.open(io.BytesIO(content)) as pdf:
        for i, page in enumerate(pdf.pages, 1):
            text = page.extract_text()
            if text and text.strip():
                lines.append(f'--- Page {i} ---')
                lines.append(text.strip())

    return '\n'.join(lines)


def parse_pptx(content: bytes, filename: str = '') -> str:
    """Extract all text from a PowerPoint presentation."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    lines = [f'[PowerPoint: {filename}]' if filename else '[PowerPoint presentation]']
    lines.append('')

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                slide_texts.append(shape.text.strip())
            # Extract table text from shapes
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_texts.append(' | '.join(cells))
        if slide_texts:
            lines.append(f'--- Slide {i} ---')
            lines.extend(slide_texts)

    return '\n'.join(lines)


def parse_xlsx(content: bytes, filename: str = '') -> str:
    """Extract all data from an Excel workbook."""
    import openpyxl

    wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines = [f'[Excel file: {filename}]' if filename else '[Excel spreadsheet]']

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_lines = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                sheet_lines.append(' | '.join(cells))
        if sheet_lines:
            lines.append(f'\n--- Sheet: {sheet_name} ---')
            lines.extend(sheet_lines)

    return '\n'.join(lines)


def parse_html(content: bytes, filename: str = '') -> str:
    """Strip HTML and extract readable text."""
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(content, 'html.parser')
    for tag in soup(['script', 'style', 'head', 'nav', 'footer']):
        tag.decompose()

    text = soup.get_text(separator='\n', strip=True)
    lines = [f'[HTML file: {filename}]' if filename else '[HTML document]']
    lines.append('')
    # Collapse excessive blank lines
    for line in text.splitlines():
        if line.strip():
            lines.append(line)

    return '\n'.join(lines)


def parse_json(content: bytes, filename: str = '') -> str:
    """Format JSON as readable structured text."""
    import json as _json

    try:
        data = _json.loads(content.decode('utf-8', errors='replace'))
        lines = [f'[JSON file: {filename}]' if filename else '[JSON data]']
        lines.append('')
        lines.append(_json.dumps(data, indent=2))
        return '\n'.join(lines)
    except _json.JSONDecodeError:
        return parse_text(content.decode('utf-8', errors='replace'))


def parse_file_by_extension(content: bytes, filename: str) -> str | None:
    """
    Route a file to the correct parser based on extension.
    Returns None for image files (handled separately via Claude vision).
    """
    ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''

    if ext == 'csv':
        return parse_csv(content.decode('utf-8', errors='replace'), filename)
    elif ext == 'eml':
        return parse_eml(content, filename)
    elif ext == 'docx':
        return parse_docx(content, filename)
    elif ext == 'pdf':
        return parse_pdf(content, filename)
    elif ext == 'pptx':
        return parse_pptx(content, filename)
    elif ext in ('xlsx',):
        return parse_xlsx(content, filename)
    elif ext in ('html', 'htm'):
        return parse_html(content, filename)
    elif ext == 'json':
        return parse_json(content, filename)
    elif ext in ('png', 'jpg', 'jpeg', 'gif', 'webp'):
        return None  # Signal to caller: needs vision analysis
    else:
        # TXT or unknown — treat as plain text
        return parse_text(content.decode('utf-8', errors='replace'))


IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
IMAGE_MIME = {
    'png':  'image/png',
    'jpg':  'image/jpeg',
    'jpeg': 'image/jpeg',
    'gif':  'image/gif',
    'webp': 'image/webp',
}


def combine_inputs(sources: list[dict]) -> str:
    """
    sources: list of {'label': str, 'text': str}
    Returns a single document with labelled sections.
    """
    parts = []
    for i, src in enumerate(sources, 1):
        label = src.get('label') or f'Source {i}'
        text = src.get('text', '').strip()
        if text:
            parts.append(f'--- {label} ---\n{text}')

    return '\n\n'.join(parts)
